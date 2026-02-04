"""PowerPoint presentation parser using python-pptx."""

from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from .base_parser import BaseParser


class PowerPointParser(BaseParser):
    """
    Parser for PowerPoint presentations (.pptx, .ppt).
    
    Extracts:
    - Slide titles and text content
    - Speaker notes
    - Shape information (images, charts, tables, groups)
    - Presentation metadata
    - AI-powered image descriptions (when enabled)
    """
    
    SUPPORTED_EXTENSIONS = ['.pptx', '.ppt']
    FILE_TYPE = 'powerpoint'
    
    def __init__(self, filepath: str, enable_ai_vision: bool = False):
        """
        Initialize PowerPoint parser.
        
        Args:
            filepath: Path to the PowerPoint file
            enable_ai_vision: Whether to enable AI image analysis
        """
        super().__init__(filepath)
        self.enable_ai_vision = enable_ai_vision
        self.ai_vision = None
        
        # Initialize AI Vision if enabled
        if self.enable_ai_vision:
            try:
                from core.ai_vision import AIVision
                self.ai_vision = AIVision()
                # Check if actually available (API key configured)
                if not self.ai_vision.is_available():
                    self.ai_vision = None
            except ImportError:
                self.ai_vision = None
    
    def parse(self) -> Dict[str, Any]:
        """
        Parse PowerPoint file and extract content from all slides.
        
        Returns:
            Dictionary containing:
            - filename, file_type, parsed_at (base fields)
            - ai_vision_enabled: Whether AI analysis was enabled
            - ai_vision_available: Whether AI analysis is actually available
            - metadata: Slide count and dimensions
            - slides: List of slide objects with content
        """
        output = self.get_base_output()
        output['ai_vision_enabled'] = self.enable_ai_vision
        output['ai_vision_available'] = self.ai_vision is not None
        
        try:
            prs = Presentation(self.filepath)
            
            # Extract metadata
            output['metadata'] = {
                'total_slides': len(prs.slides),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height
            }
            
            # Track image analysis stats
            images_analyzed = 0
            images_total = 0
            
            # Extract slides
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_data = self._parse_slide(slide, i)
                slides.append(slide_data)
                
                # Count images
                images_total += slide_data.get('image_count', 0)
                for shape in slide_data.get('shapes', []):
                    if shape.get('content_type') == 'image' and shape.get('ai_analysis', {}).get('enabled'):
                        images_analyzed += 1
            
            output['slides'] = slides
            
            # Add analysis summary
            if self.enable_ai_vision:
                output['ai_analysis_summary'] = {
                    'images_total': images_total,
                    'images_analyzed': images_analyzed,
                    'ai_available': self.ai_vision is not None
                }
            
        except Exception as e:
            output['error'] = f"PowerPoint parsing error: {str(e)}"
            output['metadata'] = {'total_slides': 0}
            output['slides'] = []
        
        return output
    
    def _parse_slide(self, slide: Slide, slide_number: int) -> Dict[str, Any]:
        """
        Parse a single slide.
        
        Args:
            slide: python-pptx Slide instance
            slide_number: 1-based slide index
            
        Returns:
            Dictionary with slide content
        """
        slide_data = {
            'slide_number': slide_number,
            'title': '',
            'text': '',
            'notes': '',
            'image_count': 0,
            'chart_count': 0,
            'table_count': 0,
            'shapes': []
        }
        
        # Extract title
        if slide.shapes.title:
            slide_data['title'] = slide.shapes.title.text.strip()
        
        # Extract speaker notes
        slide_data['notes'] = self._extract_notes(slide)
        
        # Process all shapes
        text_parts = []
        for shape in slide.shapes:
            shape_info = self._parse_shape(shape, slide_number)
            slide_data['shapes'].append(shape_info)
            
            # Collect text content
            if shape_info.get('has_text') and shape_info.get('text'):
                text_parts.append(shape_info['text'])
            
            # Update counts
            content_type = shape_info.get('content_type', '')
            if content_type == 'image':
                slide_data['image_count'] += 1
            elif content_type == 'chart':
                slide_data['chart_count'] += 1
            elif content_type == 'table':
                slide_data['table_count'] += 1
        
        # Combine all text
        slide_data['text'] = '\n'.join(text_parts)
        
        return slide_data
    
    def _extract_notes(self, slide: Slide) -> str:
        """
        Extract speaker notes from slide.
        
        Args:
            slide: python-pptx Slide instance
            
        Returns:
            Notes text or empty string
        """
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    return notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        return ''
    
    def _parse_shape(self, shape: BaseShape, slide_number: int) -> Dict[str, Any]:
        """
        Parse a single shape and determine its type and content.
        
        Args:
            shape: python-pptx shape instance
            slide_number: Current slide number for context
            
        Returns:
            Dictionary with shape information
        """
        shape_info = {
            'type': self._get_shape_type_name(shape),
            'content_type': 'unknown',
            'has_text': False,
            'text': '',
            'description': '',
            'note': ''
        }
        
        # Determine content type based on shape type
        try:
            shape_type = shape.shape_type
            
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info['content_type'] = 'image'
                shape_info['description'] = 'Image detected'
                shape_info = self._enrich_image_info(shape, shape_info, slide_number)
                
            elif shape_type == MSO_SHAPE_TYPE.CHART:
                shape_info['content_type'] = 'chart'
                shape_info['description'] = 'Chart detected'
                shape_info = self._enrich_chart_info(shape, shape_info)
                
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                shape_info['content_type'] = 'table'
                shape_info['description'] = 'Table detected'
                shape_info = self._enrich_table_info(shape, shape_info)
                
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                shape_info['content_type'] = 'group'
                shape_info['description'] = 'Grouped shapes'
                
            elif shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                shape_info['content_type'] = 'embedded_object'
                shape_info['description'] = 'Embedded OLE object (e.g., Excel chart)'
                
            elif shape_type == MSO_SHAPE_TYPE.LINKED_OLE_OBJECT:
                shape_info['content_type'] = 'linked_object'
                shape_info['description'] = 'Linked OLE object'
        
        except Exception:
            pass
        
        # Extract text if shape has text frame
        if hasattr(shape, 'text_frame'):
            try:
                text = shape.text_frame.text.strip()
                if text:
                    shape_info['has_text'] = True
                    shape_info['text'] = text
                    # Only override content_type if it's unknown
                    if shape_info['content_type'] == 'unknown':
                        shape_info['content_type'] = 'text'
            except Exception:
                pass
        
        return shape_info
    
    def _get_shape_type_name(self, shape: BaseShape) -> str:
        """
        Get human-readable shape type name.
        
        Args:
            shape: python-pptx shape instance
            
        Returns:
            Shape type name string
        """
        try:
            return str(shape.shape_type).replace('MSO_SHAPE_TYPE.', '')
        except Exception:
            return 'UNKNOWN'
    
    def _enrich_image_info(self, shape: BaseShape, shape_info: Dict, 
                           slide_number: int) -> Dict[str, Any]:
        """
        Add image-specific information to shape info.
        
        Args:
            shape: python-pptx Picture shape
            shape_info: Existing shape info dict
            slide_number: Current slide number
            
        Returns:
            Enriched shape info dict
        """
        try:
            image = shape.image
            shape_info['image_format'] = image.ext
            shape_info['image_size_bytes'] = len(image.blob)
            
            # Add dimensions if available
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                shape_info['width_emu'] = shape.width
                shape_info['height_emu'] = shape.height
            
            # AI Vision analysis
            if self.enable_ai_vision and self.ai_vision:
                try:
                    # Determine analysis type based on content hints
                    analysis_type = 'general'
                    
                    analysis = self.ai_vision.analyze_image(
                        image.blob,
                        image.ext.lower(),
                        analysis_type=analysis_type,
                        context=f"Slide {slide_number} of presentation"
                    )
                    shape_info['ai_analysis'] = analysis
                    
                    # Use AI description if successful
                    if analysis.get('enabled') and analysis.get('description'):
                        shape_info['description'] = analysis['description']
                    elif analysis.get('error'):
                        shape_info['ai_analysis_error'] = analysis['error']
                        
                except Exception as e:
                    shape_info['ai_analysis'] = {
                        'enabled': True,
                        'error': f'Analysis failed: {str(e)}'
                    }
            else:
                # Default description when AI not available
                size_kb = len(image.blob) / 1024
                shape_info['description'] = f"Image ({image.ext.upper()}, {size_kb:.1f} KB)"
            
        except Exception as e:
            shape_info['note'] = f'Could not extract image details: {str(e)}'
        
        return shape_info
    
    def _enrich_chart_info(self, shape: BaseShape, shape_info: Dict) -> Dict[str, Any]:
        """
        Add chart-specific information to shape info.
        
        Args:
            shape: python-pptx Chart shape
            shape_info: Existing shape info dict
            
        Returns:
            Enriched shape info dict
        """
        try:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                shape_info['chart_type'] = str(chart.chart_type) if hasattr(chart, 'chart_type') else 'unknown'
                shape_info['description'] = f"Chart: {shape_info['chart_type']}"
        except Exception as e:
            shape_info['note'] = f'Could not extract chart details: {str(e)}'
        
        return shape_info
    
    def _enrich_table_info(self, shape: BaseShape, shape_info: Dict) -> Dict[str, Any]:
        """
        Add table-specific information to shape info.
        
        Args:
            shape: python-pptx Table shape
            shape_info: Existing shape info dict
            
        Returns:
            Enriched shape info dict
        """
        try:
            if hasattr(shape, 'table'):
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                shape_info['table_rows'] = rows
                shape_info['table_columns'] = cols
                shape_info['description'] = f"Table: {rows} rows × {cols} columns"
                
                # Extract table data
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                shape_info['table_data'] = table_data
                
        except Exception as e:
            shape_info['note'] = f'Could not extract table details: {str(e)}'
        
        return shape_info
