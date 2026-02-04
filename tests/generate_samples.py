#!/usr/bin/env python3
"""Generate sample test files for integration testing."""

import os
import sys

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

SAMPLE_DIR = os.path.join(PROJECT_ROOT, 'tests', 'sample_files')
os.makedirs(SAMPLE_DIR, exist_ok=True)


def create_sample_word():
    """Create a sample Word document."""
    from docx import Document
    from docx.shared import Inches
    
    doc = Document()
    doc.core_properties.title = "Sample Word Document"
    doc.core_properties.author = "File Parser Agent"
    
    doc.add_heading('Sample Word Document', 0)
    doc.add_paragraph('This is a test document for the File Parser Agent.')
    
    doc.add_heading('Section 1: Introduction', level=1)
    doc.add_paragraph('This document tests the Word parser functionality.')
    
    doc.add_heading('Section 2: Sample Table', level=1)
    table = doc.add_table(rows=3, cols=3)
    table.style = 'Table Grid'
    
    # Header row
    header_cells = table.rows[0].cells
    header_cells[0].text = 'Name'
    header_cells[1].text = 'Value'
    header_cells[2].text = 'Notes'
    
    # Data rows
    data_cells = table.rows[1].cells
    data_cells[0].text = 'Item 1'
    data_cells[1].text = '100'
    data_cells[2].text = 'First item'
    
    data_cells = table.rows[2].cells
    data_cells[0].text = 'Item 2'
    data_cells[1].text = '200'
    data_cells[2].text = 'Second item'
    
    doc.add_heading('Section 3: Conclusion', level=1)
    doc.add_paragraph('The parser should extract all content from this document.')
    
    filepath = os.path.join(SAMPLE_DIR, 'sample.docx')
    doc.save(filepath)
    print(f"Created: {filepath}")
    return filepath


def create_sample_excel():
    """Create a sample Excel spreadsheet."""
    from openpyxl import Workbook
    
    wb = Workbook()
    
    # Sheet 1: Sales Data
    ws1 = wb.active
    ws1.title = "Sales Data"
    
    ws1['A1'] = 'Product'
    ws1['B1'] = 'Q1'
    ws1['C1'] = 'Q2'
    ws1['D1'] = 'Q3'
    ws1['E1'] = 'Q4'
    
    data = [
        ['Widget A', 100, 150, 200, 250],
        ['Widget B', 80, 90, 110, 130],
        ['Widget C', 200, 180, 220, 240],
    ]
    
    for i, row in enumerate(data, start=2):
        for j, value in enumerate(row):
            ws1.cell(row=i, column=j+1, value=value)
    
    # Sheet 2: Summary
    ws2 = wb.create_sheet("Summary")
    ws2['A1'] = 'Metric'
    ws2['B1'] = 'Value'
    ws2['A2'] = 'Total Products'
    ws2['B2'] = 3
    ws2['A3'] = 'Best Quarter'
    ws2['B3'] = 'Q4'
    
    filepath = os.path.join(SAMPLE_DIR, 'sample.xlsx')
    wb.save(filepath)
    print(f"Created: {filepath}")
    return filepath


def create_sample_pptx():
    """Create a sample PowerPoint presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Sample Presentation"
    subtitle.text = "Created for File Parser Agent Testing"
    
    # Content slide 1
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Features"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "PDF Parsing"
    
    p = tf.add_paragraph()
    p.text = "Word Document Parsing"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Excel Spreadsheet Parsing"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "PowerPoint Parsing"
    p.level = 0
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are the key features of the File Parser Agent."
    
    # Content slide 2 with table
    table_slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(table_slide_layout)
    
    title = slide.shapes.title
    title.text = "Data Summary"
    
    # Add a table
    rows, cols = 3, 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    table.cell(0, 0).text = "Format"
    table.cell(0, 1).text = "Extension"
    table.cell(0, 2).text = "Status"
    
    # Data rows
    table.cell(1, 0).text = "PDF"
    table.cell(1, 1).text = ".pdf"
    table.cell(1, 2).text = "Supported"
    
    table.cell(2, 0).text = "Word"
    table.cell(2, 1).text = ".docx"
    table.cell(2, 2).text = "Supported"
    
    filepath = os.path.join(SAMPLE_DIR, 'sample.pptx')
    prs.save(filepath)
    print(f"Created: {filepath}")
    return filepath


def create_sample_pdf():
    """Create a simple PDF file using reportlab if available, otherwise skip."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        
        filepath = os.path.join(SAMPLE_DIR, 'sample.pdf')
        
        c = canvas.Canvas(filepath, pagesize=letter)
        width, height = letter
        
        # Page 1
        c.setFont("Helvetica-Bold", 24)
        c.drawString(100, height - 100, "Sample PDF Document")
        
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 150, "This is a test document for the File Parser Agent.")
        c.drawString(100, height - 180, "It contains multiple pages of content.")
        
        c.drawString(100, height - 230, "Key Features:")
        c.drawString(120, height - 250, "• Text extraction from PDF files")
        c.drawString(120, height - 270, "• Page-by-page content parsing")
        c.drawString(120, height - 290, "• Metadata extraction")
        
        c.showPage()
        
        # Page 2
        c.setFont("Helvetica-Bold", 18)
        c.drawString(100, height - 100, "Page 2: Additional Content")
        
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 150, "This is the second page of the sample PDF.")
        c.drawString(100, height - 180, "The parser should extract content from all pages.")
        
        c.save()
        print(f"Created: {filepath}")
        return filepath
        
    except ImportError:
        print("Note: reportlab not installed. Skipping PDF generation.")
        print("      Install with: pip install reportlab")
        return None


if __name__ == '__main__':
    print("Generating sample test files...")
    print(f"Output directory: {SAMPLE_DIR}")
    print("-" * 50)
    
    create_sample_word()
    create_sample_excel()
    create_sample_pptx()
    create_sample_pdf()
    
    print("-" * 50)
    print("Done! Run tests/test_parsers.py to verify parsing.")
